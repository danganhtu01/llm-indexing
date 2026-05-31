"""Per-file-type text extraction, with OCR fallback for scans and images."""
from __future__ import annotations

import html as _html
import io
import re

TEXT_EXTS = {".txt", ".md", ".markdown", ".csv", ".tsv", ".log", ".json", ".xml",
             ".html", ".htm", ".yaml", ".yml", ".ini", ".cfg", ".rtf", ".srt", ".vtt"}
CODE_EXTS = {".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".c", ".h", ".cpp", ".cs",
             ".go", ".rs", ".rb", ".php", ".sql", ".sh", ".ps1", ".bat", ".r", ".css", ".scss"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp", ".gif"}
EMAIL_EXTS = {".eml", ".wdseml", ".emlx"}

CATEGORY = {
    **{e: "document" for e in (".doc", ".docx", ".odt", ".rtf", ".pages")},
    **{e: "spreadsheet" for e in (".xls", ".xlsx", ".xlsm", ".ods", ".csv", ".tsv")},
    **{e: "presentation" for e in (".ppt", ".pptx", ".odp", ".key")},
    ".pdf": "pdf",
    **{e: "email" for e in (".eml", ".wdseml", ".emlx", ".msg")},
    **{e: "image" for e in IMAGE_EXTS},
    **{e: "audio" for e in (".mp3", ".wav", ".flac", ".m4a", ".aac", ".ogg")},
    **{e: "video" for e in (".mp4", ".mkv", ".avi", ".mov", ".wmv", ".flv")},
    **{e: "archive" for e in (".zip", ".rar", ".7z", ".tar", ".gz", ".bz2")},
    **{e: "code" for e in CODE_EXTS},
    **{e: "text" for e in (".txt", ".md", ".markdown", ".log", ".srt", ".vtt")},
    **{e: "data" for e in (".json", ".xml", ".yaml", ".yml", ".db", ".sqlite", ".parquet")},
}


def category_for(ext: str) -> str:
    return CATEGORY.get(ext.lower(), "other")


def _decode(raw: bytes, max_chars: int) -> str:
    try:
        from charset_normalizer import from_bytes
        best = from_bytes(raw).best()
        if best is not None:
            return str(best)[:max_chars]
    except Exception:
        pass
    return raw.decode("utf-8", errors="ignore")[:max_chars]


def extract(path: str, ext: str, size: int, cfg, ocr) -> dict:
    """Return {'text', 'method', 'ocr_used', 'pages'} for one file."""
    ext = ext.lower()
    max_chars = cfg["max_chars"]
    ocr_mode = cfg["ocr"]
    res = {"text": "", "method": "name-only", "ocr_used": False, "pages": 0}

    if size > cfg["max_bytes"] or cfg.skip_ext(ext):
        return res
    try:
        if ext in EMAIL_EXTS:
            res["text"], res["method"] = _email(path, max_chars), "email"
        elif ext in TEXT_EXTS or ext in CODE_EXTS:
            with open(path, "rb") as f:
                res["text"] = _decode(f.read(max_chars * 4), max_chars)
            res["method"] = "text"
        elif ext == ".pdf":
            res.update(_pdf(path, cfg, ocr, ocr_mode, max_chars))
        elif ext == ".docx":
            res["text"], res["method"] = _docx(path, max_chars), "docx"
        elif ext in (".xlsx", ".xlsm"):
            res["text"], res["method"] = _xlsx(path, max_chars), "xlsx"
        elif ext == ".pptx":
            res["text"], res["method"] = _pptx(path, max_chars), "pptx"
        elif ext in IMAGE_EXTS and ocr_mode in ("auto", "on") and ocr and ocr.available:
            res["text"] = _image_ocr(path, ocr, max_chars)
            res["method"] = "ocr"
            res["ocr_used"] = bool(res["text"].strip())
    except Exception as e:
        res["method"] = f"error:{type(e).__name__}"
    return res


def _pix_to_pil(pix):
    from PIL import Image
    return Image.open(io.BytesIO(pix.tobytes("png")))


# PyMuPDF (fitz) is NOT thread-safe; concurrent use from worker threads causes
# native crashes (ACCESS_VIOLATION / "Bad file descriptor"). Serialize all of it.
import threading
_FITZ_LOCK = threading.Lock()


def _pdf(path, cfg, ocr, ocr_mode, max_chars):
    import fitz
    out = {"text": "", "method": "pdf-text", "ocr_used": False, "pages": 0}
    with _FITZ_LOCK, fitz.open(path) as doc:
        out["pages"] = doc.page_count
        parts = []
        for page in doc:
            parts.append(page.get_text("text"))
            if sum(len(p) for p in parts) > max_chars:
                break
        text = "".join(parts)
        need_ocr = ocr_mode == "on" or (
            ocr_mode == "auto" and len(text.strip()) < 20 * max(1, doc.page_count)
        )
        if need_ocr and ocr and ocr.available:
            ocr_parts = []
            for i in range(min(doc.page_count, cfg["ocr_max_pages"])):
                ocr_parts.append(ocr.image_to_text(_pix_to_pil(doc[i].get_pixmap(dpi=200))))
                if sum(len(p) for p in ocr_parts) > max_chars:
                    break
            ocr_text = "\n".join(ocr_parts)
            if len(ocr_text.strip()) > len(text.strip()):
                text = (text + "\n" + ocr_text).strip()
                out["method"], out["ocr_used"] = "pdf-ocr", True
        out["text"] = text[:max_chars]
    return out


def _docx(path, max_chars):
    import docx
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs]
    for tbl in d.tables:
        for row in tbl.rows:
            parts.append("\t".join(c.text for c in row.cells))
    return "\n".join(parts)[:max_chars]


def _xlsx(path, max_chars):
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts, total = [], 0
    try:
        for ws in wb.worksheets:
            parts.append(str(ws.title))
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    line = "\t".join(cells)
                    parts.append(line)
                    total += len(line)
                    if total > max_chars:
                        return "\n".join(parts)[:max_chars]
    finally:
        wb.close()
    return "\n".join(parts)[:max_chars]


def _pptx(path, max_chars):
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            parts.append(slide.notes_slide.notes_text_frame.text)
    return "\n".join(parts)[:max_chars]


def _image_ocr(path, ocr, max_chars):
    from PIL import Image
    with Image.open(path) as img:
        return ocr.image_to_text(img)[:max_chars]


_TAG_RE = re.compile(r"<(script|style)[^>]*>.*?</\1>", re.I | re.S)
_ANGLE_RE = re.compile(r"<[^>]+>")
_WS_RE = re.compile(r"[ \t\r\f\v]*\n\s*\n\s*", re.S)


def _strip_html(s: str) -> str:
    s = _TAG_RE.sub(" ", s)
    s = _ANGLE_RE.sub(" ", s)
    s = _html.unescape(s)
    return _WS_RE.sub("\n", s).strip()


def _email(path, max_chars):
    """Extract headers + body (plain or HTML->text) + attachment names."""
    from email import policy
    from email.parser import BytesParser
    with open(path, "rb") as f:
        msg = BytesParser(policy=policy.default).parse(f)

    parts = []
    for h in ("Subject", "From", "To", "Cc", "Date"):
        v = msg.get(h)
        if v:
            parts.append(f"{h}: {v}")

    texts = []
    try:
        body = msg.get_body(preferencelist=("plain", "html"))
    except Exception:
        body = None
    if body is not None:
        try:
            content = body.get_content()
            if body.get_content_type() == "text/html":
                content = _strip_html(content)
            texts.append(content)
        except Exception:
            pass
    if not texts:  # fallback: walk every text part
        for part in msg.walk():
            ct = part.get_content_type()
            if ct in ("text/plain", "text/html"):
                try:
                    c = part.get_content()
                    texts.append(_strip_html(c) if ct == "text/html" else c)
                except Exception:
                    pass
    parts.extend(texts)

    try:
        for att in msg.iter_attachments():
            fn = att.get_filename()
            if fn:
                parts.append(f"[attachment: {fn}]")
    except Exception:
        pass
    return "\n".join(p for p in parts if p)[:max_chars]
